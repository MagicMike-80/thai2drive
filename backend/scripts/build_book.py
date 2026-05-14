"""
Thai2Drive - Bygg laeringsbok fra PowerPoint
Leser PPTX, skriver om tekst med egne ord, oversetter til thai+engelsk, lagrer i DB.
Bruk: python build_book.py
"""

import asyncio
import json
import os
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path
from dotenv import load_dotenv
import anthropic
from motor.motor_asyncio import AsyncIOMotorClient

if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

load_dotenv(override=True)

api_key = os.environ.get("ANTHROPIC_API_KEY")
if not api_key:
    env_path = os.path.join(os.path.dirname(__file__), '..', '.env')
    with open(env_path) as f:
        for line in f:
            if line.startswith("ANTHROPIC_API_KEY="):
                api_key = line.strip().split("=", 1)[1]
                break

client = anthropic.Anthropic(api_key=api_key)

PPTX_PATH = r"C:\Users\Stein Hoang\Desktop\My Agents\4tTrafikkalt.pptx"
NAV = {'trinn', 'meny', 'tilbake', 'bytt bilde', 'neste', 'oppgaver', 'x',
       'skjul', 'vis', 'tilbake', 'tema'}

# Kapittelstruktur basert pa slide-numre
CHAPTERS = [
    {"num": 1, "title": "Grunnleggende trafikklare",    "no": "Grunnleggende trafikklære",     "slides": (1, 120)},
    {"num": 2, "title": "Mennesket i trafikken",         "no": "Mennesket i trafikken",          "slides": (121, 230)},
    {"num": 3, "title": "Trafikksystemet og regler",     "no": "Trafikksystemet og regler",      "slides": (231, 340)},
    {"num": 4, "title": "Kjoretoy og vegtrafikkloven",   "no": "Kjøretøy og vegtrafikkloven",    "slides": (341, 437)},
]


def get_slide_content(slide) -> tuple[str, str]:
    """Hent tittel og body fra en slide"""
    title = ''
    body_parts = []
    try:
        for shape in slide.shapes:
            if not hasattr(shape, 'text'):
                continue
            t = shape.text.strip()
            if not t or len(t) < 4:
                continue
            lines = [l.strip() for l in t.splitlines()
                     if l.strip() and l.strip().lower() not in NAV
                     and not l.strip().isdigit() and len(l.strip()) > 3]
            if not lines:
                continue
            if not title and len(lines[0]) < 100:
                title = lines[0]
                body_parts.extend(lines[1:])
            else:
                body_parts.extend(lines)
    except Exception:
        pass
    return title, '\n'.join(body_parts)


def extract_chapter_slides(pptx_path: str, start: int, end: int) -> list[tuple[int, str, str]]:
    """Hent alle slides i et kapittel med innhold"""
    try:
        from pptx import Presentation
    except ImportError:
        os.system(f"{sys.executable} -m pip install python-pptx -q")
        from pptx import Presentation

    prs = Presentation(pptx_path)
    result = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num < start or slide_num > end:
            continue
        try:
            title, body = get_slide_content(slide)
            full = (title + ' ' + body).strip()
            words = [w for w in full.lower().split()
                     if w not in NAV and len(w) > 3]
            if len(words) >= 15:
                result.append((slide_num, title, body))
        except Exception:
            pass

    return result


def make_sections(slides: list[tuple[int, str, str]], max_chars: int = 2500) -> list[tuple[list[int], str, str]]:
    """Grupper slides i seksjoner"""
    sections = []
    cur_slides = []
    cur_title = ''
    cur_body = ''

    for slide_num, title, body in slides:
        combined = title + '\n' + body
        if len(cur_body) + len(combined) > max_chars and cur_body:
            sections.append((cur_slides, cur_title, cur_body.strip()))
            cur_slides = [slide_num]
            cur_title = title
            cur_body = body
        else:
            cur_slides.append(slide_num)
            if not cur_title and title:
                cur_title = title
            cur_body += '\n' + combined

    if cur_body.strip():
        sections.append((cur_slides, cur_title, cur_body.strip()))

    return sections


def rewrite_and_translate(raw_text: str, section_title: str, chapter_title: str) -> dict | None:
    """Skriv om med egne ord og oversett til thai + engelsk"""

    prompt = f"""Du jobber med en laeringsbok om norsk trafikklare for kapitlet: "{chapter_title}".

Her er ratatekst fra et avsnitt (overskrift: "{section_title}"):
---
{raw_text[:2500]}
---

Gjor folgende:
1. Skriv innholdet om med DINE EGNE ORD pa naturlig, enkel norsk. Behold alle fakta korrekte. LAG NYE SETNINGER - ikke kopier fra originalen.
2. Lag en kort, beskrivende tittel for avsnittet (maks 6 ord)
3. Oversett tittelen og innholdet til thai og engelsk

Svaret skal vaere flytende tekst (ikke punktlister). 3-6 setninger per spraakvariasjon.

Svar KUN med JSON:
{{
  "tittel": {{
    "no": "Norsk tittel",
    "th": "หัวข้อภาษาไทย",
    "en": "English title"
  }},
  "innhold": {{
    "no": "Omskrevet tekst pa norsk.",
    "th": "เนื้อหาภาษาไทย",
    "en": "English content."
  }}
}}"""

    try:
        response = client.messages.create(
            model="claude-opus-4-7",
            max_tokens=1500,
            messages=[{"role": "user", "content": prompt}]
        )
        text = response.content[0].text.strip()
        if "```" in text:
            text = text.split("```")[1]
            if text.startswith("json"):
                text = text[4:]
            text = text.rsplit("```", 1)[0]
        return json.loads(text.strip())
    except Exception as e:
        print(f"Feil: {e}")
        return None


async def save_section(db, chapter: dict, section_num: int,
                       slide_nums: list[int], data: dict) -> bool:
    doc = {
        "id": str(uuid.uuid4()),
        "chapter_num": chapter["num"],
        "chapter_title": {
            "no": chapter["no"],
            "th": chapter["no"],
            "en": chapter["title"]
        },
        "section_num": section_num,
        "section_title": data["tittel"],
        "content": data["innhold"],
        "slides": slide_nums,
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.chapters.insert_one(doc)
    return True


async def main():
    start_ch = int(sys.argv[1]) if len(sys.argv) > 1 else 1
    end_ch = int(sys.argv[2]) if len(sys.argv) > 2 else 4

    mongo_url = os.environ.get("MONGO_URL")
    db_name = os.environ.get("DB_NAME", "thai2drive")
    if not mongo_url:
        env_path = os.path.join(os.path.dirname(__file__), '..', '.env')
        with open(env_path) as f:
            for line in f:
                if line.startswith("MONGO_URL="):
                    mongo_url = line.strip().split("=", 1)[1].strip('"')
                if line.startswith("DB_NAME="):
                    db_name = line.strip().split("=", 1)[1].strip('"')

    mongo_client = AsyncIOMotorClient(mongo_url)
    db = mongo_client[db_name]
    await db.chapters.create_index([("chapter_num", 1), ("section_num", 1)])

    print(f"\n{'='*60}")
    print(f"  Thai2Drive - Bygger laeringsbok (PPTX)")
    print(f"{'='*60}")
    print(f"  Fil: {Path(PPTX_PATH).name}")

    total_new = 0

    for ch in CHAPTERS:
        if not (start_ch <= ch["num"] <= end_ch):
            continue

        print(f"\n  Kapittel {ch['num']}: {ch['no']}")
        print(f"  Slide {ch['slides'][0]}-{ch['slides'][1]}")

        existing = await db.chapters.count_documents({"chapter_num": ch["num"]})
        if existing > 0:
            print(f"  Allerede ferdig ({existing} seksjoner)")
            continue

        slides = extract_chapter_slides(PPTX_PATH, ch["slides"][0], ch["slides"][1])
        print(f"  {len(slides)} slides med innhold")

        if not slides:
            print("  Ingen innhold funnet")
            continue

        sections = make_sections(slides)
        print(f"  Delt i {len(sections)} seksjoner")

        for s_num, (slide_nums, sec_title, raw_text) in enumerate(sections, 1):
            preview = sec_title[:45] if sec_title else raw_text[:45]
            print(f"  [{s_num}/{len(sections)}] {preview}...", end=" ", flush=True)
            result = rewrite_and_translate(raw_text, sec_title, ch["no"])
            if result:
                await save_section(db, ch, s_num, slide_nums, result)
                print(f"OK: {result['tittel']['no'][:40]}")
                total_new += 1
            else:
                print("Hoppet over")

    total = await db.chapters.count_documents({})
    mongo_client.close()

    print(f"\n{'='*60}")
    print(f"  FERDIG! {total_new} nye seksjoner")
    print(f"  Totalt i DB: {total} seksjoner")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    asyncio.run(main())
