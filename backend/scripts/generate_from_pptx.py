"""
generate_from_pptx.py
=====================
Leser slides fra PPTX-filer, analyserer med Claude AI,
genererer trafikkspørsmål på norsk/thai/engelsk, lagrer i MongoDB.

Kjør: python scripts/generate_from_pptx.py
"""

import asyncio
import base64
import json
import os
import sys
import uuid
from pathlib import Path
from datetime import datetime, timezone
from io import BytesIO

import anthropic
import zipfile
from motor.motor_asyncio import AsyncIOMotorClient
from dotenv import load_dotenv
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# ── Setup ──────────────────────────────────────────────────────────────
ROOT = Path(__file__).parent.parent
load_dotenv(ROOT / '.env')

PPTX_FILES = [
    (r"C:\Users\Stein Hoang\OneDrive\TGK A-Team\A-TEAM FØRSTEHJELP PP.pptx", "Safety"),
    (r"C:\Users\Stein Hoang\OneDrive\TGK A-Team\NYTT TGK A-TEAM ALLE DAGER (1).pptx", "Road Rules"),
    (r"C:\Users\Stein Hoang\CrossDevice\Michael sin Z Flip5\storage\Download\4tTrafikkalt.pptx", "Road Rules"),
    (r"C:\Users\Stein Hoang\CrossDevice\Michael sin Z Flip5\storage\Download\Undervisningsplan Fokus 2019.ppsx", "Road Rules"),
]

MONGO_URL     = os.environ['MONGO_URL']
DB_NAME       = os.environ['DB_NAME']
ANTHROPIC_KEY = os.environ.get('ANTHROPIC_API_KEY')

client_ai = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
mongo     = AsyncIOMotorClient(MONGO_URL)
db        = mongo[DB_NAME]

# ── Slide-innhold ──────────────────────────────────────────────────────
def get_slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        except:
            pass
    return "\n".join(texts)


def get_slide_image(slide, max_px=800, quality=80) -> tuple[str, str] | None:
    """Hent største bilde fra slide. Returnerer (base64, media_type) eller None."""
    best = None
    best_area = 0
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                area = shape.width * shape.height
                if area > best_area:
                    best = shape
                    best_area = area
        except:
            pass

    if not best:
        return None

    try:
        img_bytes = best.image.blob
        img = Image.open(BytesIO(img_bytes)).convert("RGB")
        w, h = img.size
        if max(w, h) > max_px:
            scale = max_px / max(w, h)
            img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=quality)
        b64 = base64.b64encode(buf.getvalue()).decode()
        return b64, "image/jpeg"
    except:
        return None


# ── AI prompt ──────────────────────────────────────────────────────────
PROMPT_TEMPLATE = """Du er ekspert på norsk trafikkopplæring og norsk trafikk­lov (lovdata.no).

Denne slide-en er fra et norsk kjørekurs. Her er innholdet:
--- SLIDE-TEKST ---
{text}
---

Lag {count} FORSKJELLIGE teoriprøvespørsmål basert på denne slide-en.
Spørsmålene skal teste forståelse, ikke bare hukommelse.

Returner KUN gyldig JSON-array:
[
  {{
    "category": "<Traffic Signs | Road Rules | Right of Way | Speed Limits | Safety | Driving Conditions | Situations | Pedestrians and Cyclists | Vehicle Knowledge | Environment and Economy>",
    "difficulty": "<easy | medium | hard>",
    "question": {{"no": "...", "th": "คำถามภาษาไทย", "en": "..."}},
    "options": [
      {{"id": "A", "text": {{"no": "...", "th": "...", "en": "..."}}}},
      {{"id": "B", "text": {{"no": "...", "th": "...", "en": "..."}}}},
      {{"id": "C", "text": {{"no": "...", "th": "...", "en": "..."}}}},
      {{"id": "D", "text": {{"no": "...", "th": "...", "en": "..."}}}}
    ],
    "correctOptionId": "A",
    "explanation": {{"no": "Forklaring med lovhenvisning", "th": "คำอธิบายภาษาไทย", "en": "Explanation with law reference"}},
    "law_reference": "Eks: Vegtrafikkloven § 3"
  }}
]

Regler:
- Kun basert på innholdet i slide-teksten
- Riktig svar iht. norsk lov / Statens vegvesen / lovdata.no
- Hvert spørsmål skal ha ULIK vinkel
- Hvis slide-teksten er for kort eller ikke egner seg → returner tom array []"""


def clean_json(text: str):
    text = text.strip()
    if "```" in text:
        parts = text.split("```")
        text = parts[1] if len(parts) > 1 else text
        if text.startswith("json"):
            text = text[4:]
    return text.strip()


async def generate_from_slide(slide_text: str, slide_img: tuple | None,
                               default_category: str, source: str,
                               slide_num: int, count: int = 2) -> list[dict]:
    if len(slide_text.strip()) < 20:
        return []  # For lite tekst

    prompt = PROMPT_TEMPLATE.format(text=slide_text, count=count)

    content = []
    if slide_img:
        b64, media_type = slide_img
        content.append({"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}})
    content.append({"type": "text", "text": prompt})

    try:
        resp = client_ai.messages.create(
            model="claude-opus-4-5",
            max_tokens=3000,
            messages=[{"role": "user", "content": content}]
        )
        text = clean_json(resp.content[0].text)
        items = json.loads(text)
        if not isinstance(items, list):
            return []

        results = []
        for data in items[:count]:
            if not data.get("question", {}).get("no"):
                continue
            # Legg til bilde hvis tilgjengelig
            if slide_img:
                b64, media_type = slide_img
                data["bildeUrl"] = f"data:{media_type};base64,{b64}"
            data["id"] = str(uuid.uuid4())
            data["schema_version"] = 2
            data["active"] = True
            data["created_at"] = datetime.now(timezone.utc).isoformat()
            data["source"] = "pptx_agent"
            data["source_file"] = source
            data["source_slide"] = slide_num
            results.append(data)
        return results

    except (json.JSONDecodeError, KeyError) as e:
        print(f"    Parse-feil: {e}")
        return []
    except Exception as e:
        print(f"    AI-feil: {e}")
        return []


async def is_duplicate(q: dict) -> bool:
    existing = await db.questions.find_one({"question.no": q["question"]["no"]})
    return existing is not None


# ── Hovedlogikk ────────────────────────────────────────────────────────
async def main():
    total_saved = 0
    total_skipped = 0
    total_errors = 0

    for pptx_path, default_cat in PPTX_FILES:
        path = Path(pptx_path)
        if not path.exists():
            print(f"IKKE FUNNET: {pptx_path}")
            continue

        print(f"\n{'='*60}")
        print(f"Fil: {path.name}")

        try:
            # PPSX-filer trenger content-type fix (in-memory)
            if path.suffix.lower() == '.ppsx':
                buf = BytesIO()
                with zipfile.ZipFile(str(path), 'r') as zin:
                    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
                        for item in zin.infolist():
                            data = zin.read(item.filename)
                            if item.filename == '[Content_Types].xml':
                                data = data.replace(
                                    b'presentationml.slideshow.main+xml',
                                    b'presentationml.presentation.main+xml'
                                )
                            zout.writestr(item, data)
                buf.seek(0)
                prs = Presentation(buf)
            else:
                prs = Presentation(str(path))
        except Exception as e:
            print(f"  Kunne ikke lese: {e}")
            continue

        slides = list(prs.slides)
        print(f"Slides: {len(slides)}")

        # Sjekk hvilke slides allerede er behandlet
        done = set()
        async for q in db.questions.find(
            {"source": "pptx_agent", "source_file": path.name},
            {"source_slide": 1}
        ):
            done.add(q.get("source_slide"))

        todo_slides = [(i, s) for i, s in enumerate(slides, 1) if i not in done]
        print(f"Allerede behandlet: {len(done)} | Gjenstaar: {len(todo_slides)}")

        for i, slide in todo_slides:
            text = get_slide_text(slide)
            img  = get_slide_image(slide)

            preview = text[:50].replace('\n', ' ') if text else "(ingen tekst)"
            print(f"  [{i}/{len(slides)}] {preview} ...", flush=True)

            if len(text.strip()) < 20:
                print(f"    HOPPET OVER (for kort tekst)")
                continue

            questions = await generate_from_slide(
                text, img, default_cat, path.name, i, count=2
            )

            if not questions:
                print(f"    INGEN spørsmål generert")
                total_errors += 1
                continue

            for q in questions:
                if await is_duplicate(q):
                    print(f"    DUPLIKAT: {q['question']['no'][:50]}")
                    total_skipped += 1
                    continue
                await db.questions.insert_one(q)
                print(f"    OK [{q.get('difficulty','?')}] {q['question']['no'][:55]}")
                total_saved += 1

            await asyncio.sleep(1.5)

    total_db = await db.questions.count_documents({})
    print(f"\n{'='*60}")
    print(f"FERDIG!")
    print(f"Lagret:   {total_saved}")
    print(f"Duplikat: {total_skipped}")
    print(f"Feil:     {total_errors}")
    print(f"Totalt i DB: {total_db} sporsmal")


if __name__ == "__main__":
    asyncio.run(main())
