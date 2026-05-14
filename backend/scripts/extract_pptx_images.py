"""
Extract images from PPTX slides and store in chapters collection.
Picks the best (largest) image from each section's slides.
Resizes to max 900px wide to keep MongoDB storage reasonable.
"""
import asyncio, os, sys, base64, io
from dotenv import load_dotenv
from motor.motor_asyncio import AsyncIOMotorClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("Tip: install Pillow for image resizing:  pip install Pillow")

if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

load_dotenv(override=True)

PPTX = r"C:\Users\Stein Hoang\Desktop\My Agents\4tTrafikkalt.pptx"
MAX_WIDTH = 900   # px
MAX_BYTES = 400_000  # 400 KB limit per image


def resize_image(blob: bytes, ext: str) -> tuple[bytes, str]:
    """Resize image to max MAX_WIDTH px wide. Returns (blob, mime_type)."""
    if not HAS_PIL:
        return blob, ext

    try:
        img = Image.open(io.BytesIO(blob))
        # Convert RGBA/P to RGB for JPEG
        if img.mode in ('RGBA', 'P', 'LA'):
            bg = Image.new('RGB', img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
            img = bg
        elif img.mode != 'RGB':
            img = img.convert('RGB')

        # Resize if too wide
        if img.width > MAX_WIDTH:
            ratio = MAX_WIDTH / img.width
            new_h = int(img.height * ratio)
            img = img.resize((MAX_WIDTH, new_h), Image.LANCZOS)

        # Save as JPEG
        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=82, optimize=True)
        result = buf.getvalue()

        # If still too large, reduce quality further
        if len(result) > MAX_BYTES:
            buf = io.BytesIO()
            img.save(buf, format='JPEG', quality=60, optimize=True)
            result = buf.getvalue()

        return result, 'jpeg'
    except Exception as e:
        print(f"    resize error: {e}")
        return blob, ext


def get_best_image(slide) -> tuple[bytes, str] | None:
    """Find the largest useful image on a slide."""
    candidates = []

    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                blob = img.blob
                ext = (img.ext or 'png').lower().replace('jpg', 'jpeg')
                # Skip tiny images (icons, bullets) and huge ones
                if 5_000 < len(blob) < 5_000_000:
                    # Prefer wider / taller shapes (more prominent)
                    area = shape.width * shape.height
                    candidates.append((area, blob, ext))
        except Exception:
            pass

    if not candidates:
        return None

    # Pick the largest by shape area
    candidates.sort(key=lambda x: x[0], reverse=True)
    _, blob, ext = candidates[0]
    return blob, ext


async def main():
    mongo_url = os.environ.get("MONGO_URL")
    db_name = os.environ.get("DB_NAME", "thai2drive")
    if not mongo_url:
        with open(os.path.join(os.path.dirname(__file__), '..', '.env')) as f:
            for line in f:
                if line.startswith("MONGO_URL="):
                    mongo_url = line.strip().split("=", 1)[1].strip('"')

    mc = AsyncIOMotorClient(mongo_url)
    db = mc[db_name]

    print("Laster PPTX...")
    prs = Presentation(PPTX)
    print(f"  {len(prs.slides)} slides lastet")

    sections = await db.chapters.find(
        {}, {"_id": 0, "id": 1, "slides": 1, "chapter_num": 1, "section_num": 1,
             "section_title": 1}
    ).sort([("chapter_num", 1), ("section_num", 1)]).to_list(300)

    print(f"  {len(sections)} seksjoner funnet i DB\n")

    updated = 0
    no_image = 0

    for sec in sections:
        slide_nums = sec.get("slides", [])
        title = sec.get("section_title", {})
        label = title.get("no") or title.get("en") or f"seksjon {sec['section_num']}"

        best = None
        for snum in slide_nums:
            if snum < 1 or snum > len(prs.slides):
                continue
            slide = prs.slides[snum - 1]
            candidate = get_best_image(slide)
            if candidate:
                blob, ext = candidate
                # Prefer the first good image we find (slides are ordered)
                if best is None or len(blob) > len(best[0]):
                    best = (blob, ext)

        if best:
            blob, ext = best
            blob, ext = resize_image(blob, ext)
            data_url = f"data:image/{ext};base64,{base64.b64encode(blob).decode()}"
            await db.chapters.update_one(
                {"id": sec["id"]},
                {"$set": {"image": data_url}}
            )
            updated += 1
            kb = len(blob) // 1024
            print(f"  ✓ {sec['chapter_num']}.{sec['section_num']} {label[:50]:<50} {kb} KB")
        else:
            no_image += 1
            print(f"  – {sec['chapter_num']}.{sec['section_num']} {label[:50]:<50} (ingen bilde)")

    print(f"\n{'='*60}")
    print(f"  Bilder lagt til: {updated}")
    print(f"  Uten bilde:      {no_image}")
    print(f"  Totalt:          {len(sections)}")

    # Show per-chapter summary
    for ch in range(1, 5):
        n = await db.chapters.count_documents({"chapter_num": ch, "image": {"$exists": True}})
        t = await db.chapters.count_documents({"chapter_num": ch})
        print(f"  Kapittel {ch}: {n}/{t} med bilde")

    mc.close()


asyncio.run(main())
