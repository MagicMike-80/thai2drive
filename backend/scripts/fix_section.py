"""Fix one failed section - usage: python fix_section.py <chapter_num> <slide_num>"""
import asyncio, json, uuid, os, sys
from datetime import datetime, timezone
from dotenv import load_dotenv
import anthropic
from motor.motor_asyncio import AsyncIOMotorClient
from pptx import Presentation

if sys.stdout.encoding != 'utf-8':
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')

load_dotenv(override=True)
client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
PPTX = r"C:\Users\Stein Hoang\Desktop\My Agents\4tTrafikkalt.pptx"
NAV = {'trinn', 'meny', 'tilbake', 'neste', 'x', 'skjul', 'vis', 'tema'}

CHAPTER_TITLES = {
    1: {"no": "Grunnleggende trafikklare",  "th": "การเรียนรู้พื้นฐานการจราจร", "en": "Basic Traffic Theory"},
    2: {"no": "Mennesket i trafikken",       "th": "มนุษย์ในการจราจร",             "en": "Humans in Traffic"},
    3: {"no": "Trafikksystemet og regler",   "th": "ระบบจราจรและกฎหมาย",           "en": "Traffic System and Rules"},
    4: {"no": "Kjoretoy og vegtrafikkloven", "th": "ยานพาหนะและกฎหมายจราจร",       "en": "Vehicles and Traffic Law"},
}

async def main():
    chapter_num = int(sys.argv[1]) if len(sys.argv) > 1 else 3
    slide_num = int(sys.argv[2]) if len(sys.argv) > 2 else 300

    mongo_url = os.environ.get("MONGO_URL")
    db_name = os.environ.get("DB_NAME", "thai2drive")
    if not mongo_url:
        with open(os.path.join(os.path.dirname(__file__), '..', '.env')) as f:
            for line in f:
                if line.startswith("MONGO_URL="):
                    mongo_url = line.strip().split("=", 1)[1].strip('"')

    mc = AsyncIOMotorClient(mongo_url)
    db = mc[db_name]

    prs = Presentation(PPTX)
    slide = prs.slides[slide_num - 1]

    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            lines = [l.strip() for l in shape.text.strip().splitlines()
                     if l.strip() and l.strip().lower() not in NAV and len(l.strip()) > 3]
            texts.extend(lines)

    raw = " ".join(texts)[:800]
    ch_title = CHAPTER_TITLES[chapter_num]
    print(f"Slide {slide_num}: {raw[:100]}...")

    prompt = f"""Skriv om dette innholdet fra norsk trafikklare med egne ord.
Skriv 3-5 setninger pa enkel norsk. Oversett til thai og engelsk. Lag en kort tittel (maks 5 ord).

Ratekst: {raw}

Svar KUN med JSON:
{{
  "tittel": {{"no": "...", "th": "...", "en": "..."}},
  "innhold": {{"no": "...", "th": "...", "en": "..."}}
}}"""

    resp = client.messages.create(
        model="claude-opus-4-7", max_tokens=1200,
        messages=[{"role": "user", "content": prompt}]
    )
    text = resp.content[0].text.strip()
    if "```" in text:
        text = text.split("```")[1]
        if text.startswith("json"):
            text = text[4:]
        text = text.rsplit("```", 1)[0]
    data = json.loads(text.strip())

    max_sec = await db.chapters.find_one({"chapter_num": chapter_num}, sort=[("section_num", -1)])
    next_num = (max_sec["section_num"] + 1) if max_sec else 1

    doc = {
        "id": str(uuid.uuid4()), "chapter_num": chapter_num,
        "chapter_title": ch_title,
        "section_num": next_num, "section_title": data["tittel"],
        "content": data["innhold"], "slides": [slide_num],
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.chapters.insert_one(doc)
    print(f"OK: {data['tittel']['no']}")

    for ch in range(1, 5):
        n = await db.chapters.count_documents({"chapter_num": ch})
        print(f"  Kapittel {ch}: {n} seksjoner")
    print(f"  Totalt: {await db.chapters.count_documents({})}")
    mc.close()

asyncio.run(main())
